from pptx import Presentation
import copy


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide = self.ppt_file.slides[slide_no]

        for shape in slide.shapes:
            print(f"slide_no:{slide_no}", shape.text)

    def copy_slide(self, from_slide_no, slide_layout_no=6):
        from_slide = self.ppt_file.slides[from_slide_no]
        # 지정한 레이아웃(6번)으로 슬라이드 추가 하기
        curr_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        for shp in from_slide.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    def save(self, filename):
        self.ppt_file.save(filename)



ppt_file_name = "재물조사라벨.pptx"

label_ppt = PowerPointAutoLabel(ppt_file_name)
for i in range(9):
    label_ppt.copy_slide(0)
label_ppt.save("[Auto]재물조사라벨.pptx")

