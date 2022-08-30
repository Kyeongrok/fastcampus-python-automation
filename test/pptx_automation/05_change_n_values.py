from pptx import Presentation
import copy


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide = self.ppt_file.slides[slide_no]

        for shape in slide.shapes:
            print(f"slide_no:{slide_no}", shape.text)

    def copy_slide(self, from_slide_no, slide_layout_no=6):
        from_slide = self.ppt_file.slides[from_slide_no]
        # 지정한 레이아웃(6번)으로 슬라이드 추가 하기
        curr_slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[slide_layout_no])

        for shp in from_slide.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    def duplicate_n_slides(self, slide_cnt, from_slide_no=0):
        for i in range(slide_cnt):
            label_ppt.copy_slide(from_slide_no)

    def get_shape_map(self, slide_no):
        slide = self.ppt_file.slides[slide_no]
        shape_map = {}
        for i, v in enumerate(slide.shapes):
            shape_map[v.name] = i
        print(f"shape_map:{shape_map}")
        return shape_map


    def change_text(self, slide_no, label_map):
        slide = self.ppt_file.slides[slide_no]
        shape_map = self.get_shape_map(slide_no)

        for shape_name, text in label_map.items():
            shape_no = shape_map[shape_name]
            slide.shapes[shape_no].text = text


    def save(self, filename):
        self.ppt_file.save(filename)


if __name__ == '__main__':

    label_ppt = PowerPointAutoLabel("재물조사라벨.pptx")
    label_ppt.duplicate_n_slides(slide_cnt=9)
    label_ppt.print_slide_shapes(0)
    label_ppt.change_text(0, {"TextBox 13":"모니터1"})
    label_ppt.change_text(1, {"TextBox 13":"모니터2"})
    label_ppt.print_slide_shapes(0)
    label_ppt.print_slide_shapes(1)
    label_ppt.save("[Auto]재물조사라벨.pptx")

