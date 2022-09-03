from pptx import Presentation


class PowerPointAutoLabel:
    ppt_file = None

    def __init__(self, filename):
        self.ppt_file = Presentation(filename)

    def print_slide_shapes(self, slide_no=0):
        slide = self.ppt_file.slides[slide_no]

        for shape in slide.shapes:
            print(f"slide_no:{slide_no}", shape.text)


if __name__ == '__main__':
    label_ppt = PowerPointAutoLabel("재물조사라벨.pptx")
    label_ppt.print_slide_shapes(0)
