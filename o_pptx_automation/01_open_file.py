from pptx import Presentation

ppt_file = Presentation("재물조사라벨.pptx")

slide0 = ppt_file.slides[0]

for shape in slide0.shapes:
    print(shape.text)
