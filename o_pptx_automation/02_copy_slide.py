from pptx import Presentation
import copy

def copy_slide_from_presentation(prs, layout_no):

    # 0번 슬라이드를 복사할 슬라이드로 지정
    ext_slide = prs.slides[0]

    # 지정한 레이아웃(6번)으로 슬라이드 추가 하기
    curr_slide = prs.slides.add_slide(prs.slide_layouts[layout_no])

    # 안에 있던 도형들 복사
    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return prs


ppt_file = Presentation("presentation1.pptx")
copy_slide_from_presentation(ppt_file, 6)

ppt_file.save("copied.pptx")