from pptx import Presentation
import copy

def copy_slide_from_presentation(prs, layout_no):

    # 0번 슬라이드를 복사할 슬라이드로 지정
    ext_slide = prs.slides[0]

    # 지정한 레이아웃(6번)으로 슬라이드 추가 하기
    curr_slide = prs.slides.add_slide(prs.slide_layouts[layout_no])

    # 안에 있던 도형들 복사
    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return prs


def change_text(ppt_file, slide_no, shape_no, target_text):
    slide = ppt_file.slides[slide_no]
    slide.shapes[shape_no].text = target_text


ppt_file = Presentation("재물조사라벨.pptx")
copy_slide_from_presentation(ppt_file, 6)
change_text(ppt_file, 0, shape_no=0, target_text="hello")
change_text(ppt_file, 0, shape_no=1, target_text="bye")

ppt_file.save("copied.pptx")