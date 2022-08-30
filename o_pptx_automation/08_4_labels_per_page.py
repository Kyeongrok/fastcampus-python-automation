from pptx import Presentation
import copy, os
import pandas as pd
from pptx.dml.color import RGBColor
import openpyxl

def copy_slide_from_presentation(prs, layout_no):

    # 0번 슬라이드를 복사할 슬라이드로 지정
    ext_slide = prs.slides[0]

    # 지정한 레이아웃(6번)으로 슬라이드 추가 하기
    curr_slide = prs.slides.add_slide(prs.slide_layouts[layout_no])

    img_dict = {}
    # 안에 있던 도형들 복사
    for shp in ext_slide.shapes:
        print(shp.name)
        if 'Picture' in shp.name:
            # image저장
            with open(shp.name + '.jpg', 'wb') as f:
                f.write(shp.image.blob)

            # image를 dict에 추가
            img_dict[shp.name + '.jpg'] = [shp.left, shp.top, shp.width, shp.height]
        else:
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    print(img_dict)
    for k, v in img_dict.items():
        curr_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)
    return prs


def change_text(ppt_file, slide_no, object_name, target_text):
    slide = ppt_file.slides[slide_no]
    shape_idx = {}
    for i, v in enumerate(slide.shapes):
        shape_idx[v.name] = i
    print(shape_idx, f'product_no_1 {shape_idx["product_no_1"]}')

    shape_no = shape_idx[object_name]
    slide.shapes[shape_no].text = target_text
    for paragraph in slide.shapes[shape_no].text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 0)


ppt_file = Presentation("재물조사라벨2.pptx")
df = pd.read_excel("품명_모델NO.xlsx" )
print(df['품명'].count())
for _ in range(int(df['품명'].count()) // 4):
    ppt_file = copy_slide_from_presentation(ppt_file, 6)

for i, row in df.iterrows():
    page = i // 4
    model_no_idx = (i - page * 4) + 1
    change_text(ppt_file, page, object_name=f'product_no_{model_no_idx}', target_text=row['모델No'])
    change_text(ppt_file, page, object_name=f'product_name_{model_no_idx}', target_text=row['품명'])
    print(page, model_no_idx)
ppt_file.save("n_per_page.pptx")
