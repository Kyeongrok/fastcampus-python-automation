from pptx import Presentation
import copy, os
import pandas as pd
from pptx.dml.color import RGBColor
import openpyxl

def copy_slide_from_presentation(prs, layout_no):

    # 0번 슬라이드를 복사할 슬라이드로 지정
    ext_slide = prs.slides[0]

    # 지정한 레이아웃(6번)으로 슬라이드 추가 하기
    curr_slide = prs.slides.add_slide(prs.slide_layouts[layout_no])

    img_dict = {}
    # 안에 있던 도형들 복사
    for shp in ext_slide.shapes:
        if 'Picture' in shp.name:
            # image저장
            with open(shp.name + '.jpg', 'wb') as f:
                f.write(shp.image.blob)

            # image를 dict에 추가
            img_dict[shp.name + '.jpg'] = [shp.left, shp.top, shp.width, shp.height]
        else:
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for k, v in img_dict.items():
        curr_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)
    return prs


def change_text(ppt_file, slide_no, shape_no, target_text):
    slide = ppt_file.slides[slide_no]
    shape_idx = {}
    for i, v in enumerate(slide.shapes):
        shape_idx[v.name] = i
    print(shape_idx)
    slide.shapes[shape_no].text = target_text
    for paragraph in slide.shapes[shape_no].text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(0, 0, 0)


def set_label(ppt_file, slide_no, product_name, product_no):
    ppt_file = copy_slide_from_presentation(ppt_file, 6)
    change_text(ppt_file, slide_no, shape_no=5, target_text=product_name)
    change_text(ppt_file, slide_no, shape_no=6, target_text=product_no)

ppt_file = Presentation("재물조사라벨2.pptx")
df = pd.read_excel("품명_모델NO.xlsx" )
print(df['품명'].count())
page = 0
change_text(ppt_file, page, shape_no=0, target_text='품명1')
# change_text(ppt_file, page, shape_no=1, target_text='모델No1')
#
# change_text(ppt_file, page, shape_no=7, target_text='품명2')
# change_text(ppt_file, page, shape_no=8, target_text='모델No2')
#
# change_text(ppt_file, page, shape_no=14, target_text='품명3')
# change_text(ppt_file, page, shape_no=15, target_text='모델No3')
#
# change_text(ppt_file, page, shape_no=21, target_text='품명4')
# change_text(ppt_file, page, shape_no=22, target_text='모델No4')
ppt_file.save("n_per_page.pptx")
