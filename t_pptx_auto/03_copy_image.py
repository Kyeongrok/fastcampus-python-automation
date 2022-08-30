# python-pptx
from pptx import Presentation
import copy, os

def copy_slide():
    ppt_file = Presentation("재물조사라벨.pptx")
    slide = ppt_file.slides[0]

    # 6 빈 레이아웃
    target_slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])

    img_dict = {}
    # 도형 복사
    for shape in slide.shapes:
        if 'Picture' in shape.name:
            with open(shape.name + '.jpg', 'wb') as f:
                f.write(shape.image.blob)
            img_dict[shape.name + '.jpg'] = [shape.left, shape.top, shape.width, shape.height]
        else:
            el = shape.element
            new_element = copy.deepcopy(el)
            target_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
        # print(shape.text)

    for k, v in img_dict.items():
        target_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

    ppt_file.save("add_slide.pptx")

copy_slide()


