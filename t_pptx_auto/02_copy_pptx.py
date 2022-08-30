# python-pptx
from pptx import Presentation
import copy

def copy_slide():

    ppt_file = Presentation("재물조사라벨.pptx")
    slide = ppt_file.slides[0]

    # 6 빈 레이아웃
    target_slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])

    # 도형 복사
    for shape in slide.shapes:
        el = shape.element
        new_element = copy.deepcopy(el)
        target_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
        # print(shape.text)

    ppt_file.save("add_slide.pptx")

copy_slide()


